import PyPDF2
import docx2txt
import pandas as pd
import pytesseract
from PIL import Image
from pptx import Presentation

class DataExtractor:
    def __init__(self, config):
        self.techniques = config['techniques']

    def extract_data(self, file_path):
        file_extension = file_path.split('.')[-1].lower()
        
        if file_extension == 'pdf':
            return self.extract_from_pdf(file_path)
        elif file_extension == 'docx':
            return self.extract_from_docx(file_path)
        elif file_extension == 'xlsx':
            return self.extract_from_xlsx(file_path)
        elif file_extension == 'txt':
            return self.extract_from_txt(file_path)
        elif file_extension == 'pptx':
            return self.extract_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

    def extract_from_pdf(self, file_path):
        text = ""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text()
            
            if not text.strip() and 'ocr' in self.techniques:
                text = self.perform_ocr(file_path)
            
            return text
        except Exception as e:
            print(f"Error extracting text from PDF: {str(e)}")
            return ""

    def extract_from_docx(self, file_path):
        return docx2txt.process(file_path)

    def extract_from_xlsx(self, file_path):
        df = pd.read_excel(file_path)
        return df.to_string()

    def extract_from_txt(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def extract_from_pptx(self, file_path):
        text = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        return '\n'.join(text)

    def perform_ocr(self, file_path):
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)